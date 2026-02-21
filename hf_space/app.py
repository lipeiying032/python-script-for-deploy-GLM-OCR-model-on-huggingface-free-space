#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GLM-OCR OpenAI Compatible API Server
HuggingFace Space 免费部署版
支持 Chatbox 等客户端直接接入
作者: GLM-OCR Deploy Script
"""

import os
import io
import sys
import json
import time
import base64
import traceback
import mimetypes
import zipfile
from pathlib import Path
from typing import Optional, List, Union
from contextlib import asynccontextmanager

from fastapi import FastAPI, HTTPException, Depends, Request
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel
import uvicorn
from PIL import Image
import requests

# ─────────────────────────── 配置 ─────────────────────────────────────────────
MODEL_NAME = "zai-org/GLM-OCR"
MODEL_ALIAS = "glm-ocr"
API_KEY = os.environ.get("API_KEY", "")   # 从 HF Space Secrets 读取
PORT = int(os.environ.get("PORT", 7860))

print(f"[STARTUP] GLM-OCR API Server v1.0")
print(f"[STARTUP] Model: {MODEL_NAME}")
print(f"[STARTUP] Port: {PORT}")
print(f"[STARTUP] API Key protection: {'ENABLED' if API_KEY else 'DISABLED (set API_KEY secret!)'}")

# ─────────────────────────── 全局模型 ─────────────────────────────────────────
_processor = None
_model = None

def load_model():
    global _processor, _model
    try:
        print("[MODEL] Loading transformers...")
        import torch
        from transformers import AutoProcessor, AutoModelForImageTextToText

        print("[MODEL] Downloading/Loading AutoProcessor...")
        _processor = AutoProcessor.from_pretrained(
            MODEL_NAME,
            trust_remote_code=True
        )

        print("[MODEL] Downloading/Loading AutoModelForImageTextToText...")
        _model = AutoModelForImageTextToText.from_pretrained(
            pretrained_model_name_or_path=MODEL_NAME,
            torch_dtype="auto",
            device_map="auto",
            trust_remote_code=True,
        )
        device = next(_model.parameters()).device
        print(f"[MODEL] Model loaded OK on device: {device}")
    except Exception:
        print("[MODEL][FATAL] Failed to load model:")
        traceback.print_exc()
        sys.exit(1)

@asynccontextmanager
async def lifespan(app: FastAPI):
    load_model()
    yield

# ─────────────────────────── FastAPI ──────────────────────────────────────────
app = FastAPI(
    title="GLM-OCR OpenAI Compatible API",
    version="1.0.0",
    lifespan=lifespan,
)
security = HTTPBearer(auto_error=False)

# ─────────────────────────── 鉴权 ─────────────────────────────────────────────
def verify_api_key(credentials: Optional[HTTPAuthorizationCredentials] = Depends(security)):
    if not API_KEY:
        return True  # 未配置 secret 时跳过
    if credentials is None:
        raise HTTPException(
            status_code=401,
            detail="Missing API Key. Add header: Authorization: Bearer YOUR_API_KEY"
        )
    if credentials.credentials != API_KEY:
        raise HTTPException(status_code=401, detail="Invalid API Key")
    return True

# ─────────────────────────── Pydantic 数据模型 ────────────────────────────────
class ImageUrlObj(BaseModel):
    url: str
    detail: Optional[str] = "auto"

class ContentPart(BaseModel):
    type: str
    text: Optional[str] = None
    image_url: Optional[ImageUrlObj] = None

class Message(BaseModel):
    role: str
    content: Union[str, List[ContentPart]]

class ChatRequest(BaseModel):
    model: Optional[str] = MODEL_ALIAS
    messages: List[Message]
    max_tokens: Optional[int] = 8192
    temperature: Optional[float] = 0.1
    stream: Optional[bool] = False

# ─────────────────────────── 文件处理工具 ─────────────────────────────────────

def b64_to_image(data_uri: str) -> Image.Image:
    """base64 data URI → PIL Image"""
    try:
        data = data_uri.split(",", 1)[1] if "," in data_uri else data_uri
        return Image.open(io.BytesIO(base64.b64decode(data))).convert("RGB")
    except Exception:
        print("[FILE][ERROR] base64 decode failed:")
        traceback.print_exc()
        raise

def url_to_image(url: str) -> Image.Image:
    """URL → PIL Image"""
    try:
        print(f"[FILE] Downloading image: {url[:80]}")
        r = requests.get(url, timeout=30, headers={"User-Agent": "GLM-OCR/1.0"})
        r.raise_for_status()
        return Image.open(io.BytesIO(r.content)).convert("RGB")
    except Exception:
        print("[FILE][ERROR] URL image download failed:")
        traceback.print_exc()
        raise

def pdf_to_images(pdf_bytes: bytes) -> List[Image.Image]:
    """PDF → List[PIL Image]"""
    try:
        from pdf2image import convert_from_bytes
        imgs = convert_from_bytes(pdf_bytes, dpi=150)
        print(f"[FILE] PDF converted: {len(imgs)} pages")
        return imgs
    except ImportError:
        print("[FILE][WARN] pdf2image not installed, skipping PDF")
        return []
    except Exception:
        print("[FILE][ERROR] PDF processing failed:")
        traceback.print_exc()
        return []

def docx_to_content(docx_bytes: bytes):
    """DOCX → (text_str, [PIL Image])"""
    try:
        import docx as python_docx
        doc = python_docx.Document(io.BytesIO(docx_bytes))
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    blob = rel.target_part.blob
                    images.append(Image.open(io.BytesIO(blob)).convert("RGB"))
                except Exception:
                    pass
        return "\n".join(texts), images
    except ImportError:
        print("[FILE][WARN] python-docx not installed")
        return "", []
    except Exception:
        print("[FILE][ERROR] DOCX processing failed:")
        traceback.print_exc()
        return "", []

def xlsx_to_text(xlsx_bytes: bytes) -> str:
    """XLSX → plain text table"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), read_only=True)
        lines = []
        for name in wb.sheetnames:
            lines.append(f"=== Sheet: {name} ===")
            for row in wb[name].iter_rows(values_only=True):
                row_str = "\t".join("" if c is None else str(c) for c in row)
                if row_str.strip():
                    lines.append(row_str)
        return "\n".join(lines)
    except ImportError:
        print("[FILE][WARN] openpyxl not installed")
        return ""
    except Exception:
        print("[FILE][ERROR] XLSX processing failed:")
        traceback.print_exc()
        return ""

def pptx_to_text(pptx_bytes: bytes) -> str:
    """PPTX → plain text"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)
    except ImportError:
        print("[FILE][WARN] python-pptx not installed")
        return ""
    except Exception:
        print("[FILE][ERROR] PPTX processing failed:")
        traceback.print_exc()
        return ""

def zip_to_text(zip_bytes: bytes) -> str:
    """ZIP → extract text from supported files inside"""
    try:
        parts = []
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            for name in zf.namelist():
                ext = Path(name).suffix.lower()
                try:
                    data = zf.read(name)
                    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"):
                        parts.append(f"[{name}]\n{data.decode('utf-8', errors='replace')}")
                    elif ext == ".xlsx":
                        parts.append(f"[{name}]\n{xlsx_to_text(data)}")
                    elif ext == ".pptx":
                        parts.append(f"[{name}]\n{pptx_to_text(data)}")
                    elif ext == ".docx":
                        text, _ = docx_to_content(data)
                        parts.append(f"[{name}]\n{text}")
                except Exception as e:
                    print(f"[FILE][WARN] ZIP entry {name} failed: {e}")
        return "\n\n".join(parts)
    except Exception:
        print("[FILE][ERROR] ZIP processing failed:")
        traceback.print_exc()
        return ""

def url_bytes(url: str):
    """URL → (bytes, ext)"""
    try:
        r = requests.get(url, timeout=30, headers={"User-Agent": "GLM-OCR/1.0"})
        r.raise_for_status()
        ct = r.headers.get("Content-Type", "")
        ext = mimetypes.guess_extension(ct.split(";")[0].strip()) or \
              Path(url.split("?")[0]).suffix.lower()
        return r.content, ext.lower()
    except Exception:
        print(f"[FILE][ERROR] URL download failed: {url}")
        traceback.print_exc()
        return None, ""

# ─────────────────────────── GLM-OCR 推理 ─────────────────────────────────────

def glm_ocr_infer(images: List[Image.Image], prompt: str = "Text Recognition:") -> str:
    """对图片列表执行 GLM-OCR 推理，返回合并文本"""
    import torch
    if not images:
        return ""
    results = []
    for idx, img in enumerate(images):
        print(f"[OCR] Inferring image {idx+1}/{len(images)} ...")
        try:
            messages = [{
                "role": "user",
                "content": [
                    {"type": "image", "image": img},
                    {"type": "text", "text": prompt},
                ],
            }]
            inputs = _processor.apply_chat_template(
                messages,
                tokenize=True,
                add_generation_prompt=True,
                return_dict=True,
                return_tensors="pt",
            ).to(_model.device)
            inputs.pop("token_type_ids", None)

            with torch.no_grad():
                gen_ids = _model.generate(**inputs, max_new_tokens=8192, do_sample=False)

            output = _processor.decode(
                gen_ids[0][inputs["input_ids"].shape[1]:],
                skip_special_tokens=True,
            ).strip()
            print(f"[OCR] Image {idx+1} done, {len(output)} chars")
            results.append(output)
        except Exception:
            print(f"[OCR][ERROR] Inference failed on image {idx+1}:")
            traceback.print_exc()
            results.append("")
    return "\n\n---\n\n".join(results)

# ─────────────────────────── 消息解析 ─────────────────────────────────────────

def parse_messages(messages: List[Message]):
    """从 OpenAI 消息列表提取: images列表 + text_prompt"""
    images = []
    text_parts = []
    ocr_instruction = "Text Recognition:"  # 默认 OCR 指令

    for msg in messages:
        if msg.role not in ("user", "system"):
            continue
        content = msg.content
        if isinstance(content, str):
            text_parts.append(content)
            continue
        for part in content:
            if part.type == "text" and part.text:
                text_parts.append(part.text)
            elif part.type == "image_url" and part.image_url:
                url_val = part.image_url.url
                try:
                    if url_val.startswith("data:"):
                        # base64 内联图片
                        images.append(b64_to_image(url_val))
                    elif any(url_val.lower().endswith(ext) for ext in
                             (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp")):
                        images.append(url_to_image(url_val))
                    else:
                        # 通用 URL：下载后判断类型
                        data, ext = url_bytes(url_val)
                        if data:
                            if ext in (".pdf",):
                                imgs = pdf_to_images(data)
                                images.extend(imgs)
                            elif ext in (".docx", ".doc"):
                                txt, imgs = docx_to_content(data)
                                if txt:
                                    text_parts.append(txt)
                                images.extend(imgs)
                            elif ext in (".xlsx", ".xls"):
                                text_parts.append(xlsx_to_text(data))
                            elif ext in (".pptx", ".ppt"):
                                text_parts.append(pptx_to_text(data))
                            elif ext in (".zip",):
                                text_parts.append(zip_to_text(data))
                            elif ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm"):
                                text_parts.append(data.decode("utf-8", errors="replace"))
                            else:
                                # 尝试当图片处理
                                try:
                                    images.append(Image.open(io.BytesIO(data)).convert("RGB"))
                                except Exception:
                                    print(f"[WARN] Unknown file type: {ext}, skipping")
                except Exception:
                    print(f"[ERROR] Failed to process content part:")
                    traceback.print_exc()

    combined_text = "\n".join(text_parts).strip()
    if combined_text:
        ocr_instruction = combined_text
    return images, ocr_instruction

# ─────────────────────────── API 端点 ─────────────────────────────────────────

@app.get("/")
def root():
    return {
        "service": "GLM-OCR OpenAI Compatible API",
        "model": MODEL_ALIAS,
        "status": "running",
        "endpoints": {
            "models": "GET /v1/models",
            "chat": "POST /v1/chat/completions",
        },
        "chatbox_config": {
            "api_url": "https://YOUR_USERNAME-YOUR_SPACE_NAME.hf.space",
            "model": MODEL_ALIAS,
            "note": "Set API_KEY in HF Space Secrets"
        }
    }

@app.get("/v1/models", dependencies=[Depends(verify_api_key)])
def list_models():
    return {
        "object": "list",
        "data": [{
            "id": MODEL_ALIAS,
            "object": "model",
            "created": int(time.time()),
            "owned_by": "zai-org",
            "permission": [],
            "root": MODEL_ALIAS,
        }]
    }

@app.post("/v1/chat/completions", dependencies=[Depends(verify_api_key)])
async def chat_completions(req: ChatRequest):
    start_time = time.time()
    request_id = f"chatcmpl-{int(start_time * 1000)}"
    print(f"\n[REQUEST] {request_id} | model={req.model} | stream={req.stream}")

    try:
        images, prompt = parse_messages(req.messages)
        print(f"[REQUEST] images={len(images)} | prompt_len={len(prompt)}")

        if images:
            # 有图片，运行 OCR
            result_text = glm_ocr_infer(images, prompt)
            if not result_text.strip():
                result_text = "(OCR returned empty result)"
        elif prompt.strip():
            # 纯文本：直接用 glm-ocr 做问答
            images_empty = []
            result_text = glm_ocr_infer(images_empty, prompt)
            if not result_text:
                result_text = "Please provide an image or document for OCR processing."
        else:
            result_text = "Please send an image or document to process."

        elapsed = time.time() - start_time
        print(f"[REQUEST] {request_id} done in {elapsed:.1f}s | result_len={len(result_text)}")

        response_obj = {
            "id": request_id,
            "object": "chat.completion",
            "created": int(start_time),
            "model": MODEL_ALIAS,
            "choices": [{
                "index": 0,
                "message": {
                    "role": "assistant",
                    "content": result_text,
                },
                "finish_reason": "stop",
            }],
            "usage": {
                "prompt_tokens": len(prompt.split()),
                "completion_tokens": len(result_text.split()),
                "total_tokens": len(prompt.split()) + len(result_text.split()),
            }
        }

        if req.stream:
            # SSE streaming (单块发出)
            def event_stream():
                chunk = {
                    "id": request_id,
                    "object": "chat.completion.chunk",
                    "created": int(start_time),
                    "model": MODEL_ALIAS,
                    "choices": [{
                        "index": 0,
                        "delta": {"role": "assistant", "content": result_text},
                        "finish_reason": None,
                    }]
                }
                yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
                # 发送结束标志
                end_chunk = {
                    "id": request_id,
                    "object": "chat.completion.chunk",
                    "created": int(start_time),
                    "model": MODEL_ALIAS,
                    "choices": [{
                        "index": 0,
                        "delta": {},
                        "finish_reason": "stop",
                    }]
                }
                yield f"data: {json.dumps(end_chunk)}\n\n"
                yield "data: [DONE]\n\n"
            return StreamingResponse(event_stream(), media_type="text/event-stream")

        return JSONResponse(content=response_obj)

    except HTTPException:
        raise
    except Exception:
        print(f"[REQUEST][ERROR] {request_id} unhandled exception:")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=traceback.format_exc())

# ─────────────────────────── 启动 ─────────────────────────────────────────────
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=PORT, log_level="info")
